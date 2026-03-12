"""pptx_inspector.py — Analyze a .pptx template to extract theme, colors, fonts, and layouts.

Produces a structured report that Claude uses to plan scenario generation.

Usage:
    python3 pptx_inspector.py template.pptx
    python3 pptx_inspector.py template.pptx --json
"""

from __future__ import annotations

import json
import os
import sys
from collections import Counter
from dataclasses import dataclass, field, asdict

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class ThemeColor:
    role: str        # e.g. "dk1", "accent1", "hlink"
    hex: str         # e.g. "#054E5A"
    source: str      # "srgb" or "system"

@dataclass
class FontScheme:
    name: str
    major_latin: str  # Heading font
    major_ea: str
    minor_latin: str  # Body font
    minor_ea: str

@dataclass
class LayoutPlaceholder:
    idx: int
    name: str
    width: int   # Emu
    height: int  # Emu

@dataclass
class SlideLayout:
    index: int
    name: str
    placeholders: list[LayoutPlaceholder] = field(default_factory=list)
    has_title: bool = False
    has_subtitle: bool = False
    has_content: bool = False
    has_picture: bool = False
    has_footer: bool = False
    is_dark: bool = False

@dataclass
class TemplateReport:
    file_path: str
    file_size_kb: float
    theme_name: str
    color_scheme_name: str
    colors: list[ThemeColor] = field(default_factory=list)
    fonts: FontScheme | None = None
    slide_width: int = 0   # Emu
    slide_height: int = 0  # Emu
    aspect_ratio: str = ""
    layouts: list[SlideLayout] = field(default_factory=list)
    existing_slides: int = 0
    recommended_colors: dict = field(default_factory=dict)
    recommended_font: str = ""
    layout_mapping: dict = field(default_factory=dict)
    is_custom_theme: bool = False       # True if theme has non-default colors
    design_colors: list[tuple] = field(default_factory=list)  # (hex, count) from shapes


# ---------------------------------------------------------------------------
# Color role mapping — theme roles to presentation semantics
# ---------------------------------------------------------------------------

ROLE_MAP = {
    "dk2":     "primary",
    "accent1": "primary_alt",
    "accent3": "accent",
    "accent5": "secondary",
    "accent4": "sage",
    "accent6": "coral",
    "accent2": "gray",
    "lt2":     "gray_light",
    "dk1":     "text",
    "lt1":     "white",
    "hlink":   "link",
    "folHlink": "link_visited",
}


# ---------------------------------------------------------------------------
# Inspector
# ---------------------------------------------------------------------------

def inspect_template(path: str) -> TemplateReport:
    """Analyze a .pptx template and return a structured report."""
    prs = Presentation(path)
    report = TemplateReport(
        file_path=os.path.abspath(path),
        file_size_kb=os.path.getsize(path) / 1024,
        theme_name="",
        color_scheme_name="",
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        existing_slides=len(prs.slides),
    )

    # Aspect ratio
    w, h = prs.slide_width, prs.slide_height
    ratio = w / h
    if abs(ratio - 16/9) < 0.01:
        report.aspect_ratio = "16:9"
    elif abs(ratio - 4/3) < 0.01:
        report.aspect_ratio = "4:3"
    else:
        report.aspect_ratio = f"{ratio:.2f}:1"

    # --- Theme extraction ---
    _extract_theme(prs, report)

    # --- Shape color extraction (for Canva/Slidesgo-style templates) ---
    _extract_design_colors(prs, report)

    # --- Layout cataloging ---
    _extract_layouts(prs, report)

    # --- Recommendations ---
    _build_recommendations(report)

    return report


def _extract_theme(prs: Presentation, report: TemplateReport):
    """Extract theme colors and fonts from slide master."""
    master = prs.slide_masters[0]
    theme_part = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            theme_part = rel.target_part
            break

    if not theme_part:
        return

    theme_xml = etree.fromstring(theme_part.blob)
    report.theme_name = theme_xml.get("name", "unknown")

    # Colors
    for cs in theme_xml.iter(qn("a:clrScheme")):
        report.color_scheme_name = cs.get("name", "unknown")
        for child in cs:
            role = child.tag.split("}")[-1]
            srgb = child.find(qn("a:srgbClr"))
            sys_clr = child.find(qn("a:sysClr"))
            if srgb is not None:
                report.colors.append(ThemeColor(
                    role=role, hex=f"#{srgb.get('val')}", source="srgb"
                ))
            elif sys_clr is not None:
                last = sys_clr.get("lastClr", sys_clr.get("val", "000000"))
                report.colors.append(ThemeColor(
                    role=role, hex=f"#{last}", source="system"
                ))

    # Fonts
    for fs in theme_xml.iter(qn("a:fontScheme")):
        major = fs.find(qn("a:majorFont"))
        minor = fs.find(qn("a:minorFont"))
        report.fonts = FontScheme(
            name=fs.get("name", "unknown"),
            major_latin=_get_typeface(major, "latin"),
            major_ea=_get_typeface(major, "ea"),
            minor_latin=_get_typeface(minor, "latin"),
            minor_ea=_get_typeface(minor, "ea"),
        )


def _luminance(hex_color: str) -> float:
    """Relative luminance of a hex color (0=black, 1=white)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
    return 0.299 * r + 0.587 * g + 0.114 * b


def _get_typeface(parent, tag: str) -> str:
    if parent is None:
        return ""
    el = parent.find(qn(f"a:{tag}"))
    return el.get("typeface", "") if el is not None else ""


def _extract_design_colors(prs: Presentation, report: TemplateReport):
    """Extract actual colors used in slide shapes.

    Many templates (Canva, Slidesgo, etc.) use default Office theme
    but embed their real design colors directly in shapes. This function
    collects those colors and detects whether the theme is custom or default.
    """
    # Check if theme is default Office
    DEFAULT_OFFICE_SCHEMES = {"Office", "Office Theme"}
    report.is_custom_theme = report.color_scheme_name not in DEFAULT_OFFICE_SCHEMES

    # Collect colors from slide shapes
    color_counter = Counter()
    SKIP_COLORS = {"000000", "FFFFFF", "ffffff"}  # Black/white are universal
    for slide in prs.slides:
        for shape in slide.shapes:
            for srgb in shape._element.iter(qn("a:srgbClr")):
                val = srgb.get("val", "")
                if val and val.upper() not in {"000000", "FFFFFF"}:
                    color_counter[f"#{val}"] += 1

    # Store top colors (at least 2 occurrences to filter noise)
    report.design_colors = [
        (hex_val, count) for hex_val, count in color_counter.most_common(12)
        if count >= 2
    ]


def _extract_layouts(prs: Presentation, report: TemplateReport):
    """Catalog all slide layouts and their placeholders."""
    for i, layout in enumerate(prs.slide_layouts):
        sl = SlideLayout(index=i, name=layout.name)
        sl.is_dark = "dark" in layout.name.lower()

        for ph in layout.placeholders:
            p = LayoutPlaceholder(
                idx=ph.placeholder_format.idx,
                name=ph.name,
                width=ph.width,
                height=ph.height,
            )
            sl.placeholders.append(p)
            name_lower = ph.name.lower()
            if "title" in name_lower:
                sl.has_title = True
            if "subtitle" in name_lower:
                sl.has_subtitle = True
            if "content" in name_lower:
                sl.has_content = True
            if "picture" in name_lower:
                sl.has_picture = True
            if "footer" in name_lower:
                sl.has_footer = True

        report.layouts.append(sl)


def _build_recommendations(report: TemplateReport):
    """Build recommended color map and layout mapping."""
    if report.is_custom_theme:
        # Custom theme — use theme colors
        for tc in report.colors:
            semantic = ROLE_MAP.get(tc.role)
            if semantic:
                report.recommended_colors[semantic] = tc.hex
    else:
        # Default Office theme — use design colors from shapes instead
        if report.design_colors:
            design = [hex_val for hex_val, _ in report.design_colors]
            # Classify colors by luminance and assign semantically
            dark_colors = []
            mid_colors = []
            light_colors = []
            for h in design:
                lum = _luminance(h)
                if lum < 0.25:
                    dark_colors.append(h)
                elif lum > 0.75:
                    light_colors.append(h)
                else:
                    mid_colors.append(h)

            # Assign semantically: primary=most-used vibrant, text=darkest, bg=lightest
            all_ordered = design  # frequency order

            # Text: darkest color
            if dark_colors:
                report.recommended_colors["text"] = dark_colors[0]

            # Primary: prefer mid-tone vibrant color, fallback to most-used non-text
            if mid_colors:
                report.recommended_colors["primary"] = mid_colors[0]
            else:
                for h in all_ordered:
                    if h != report.recommended_colors.get("text"):
                        report.recommended_colors["primary"] = h
                        break
            if "primary" not in report.recommended_colors:
                report.recommended_colors["primary"] = all_ordered[0]

            # Gray bg: lightest color (not already used as primary)
            for h in light_colors:
                if h != report.recommended_colors.get("primary"):
                    report.recommended_colors["gray_bg"] = h
                    break

            # Fill remaining semantic roles from unused colors
            used = set(report.recommended_colors.values())
            EXTRA_ROLES = ["accent", "secondary", "sage", "coral"]
            role_idx = 0
            for h in all_ordered:
                if h not in used and role_idx < len(EXTRA_ROLES):
                    report.recommended_colors[EXTRA_ROLES[role_idx]] = h
                    used.add(h)
                    role_idx += 1
        else:
            # Fallback to theme
            for tc in report.colors:
                semantic = ROLE_MAP.get(tc.role)
                if semantic:
                    report.recommended_colors[semantic] = tc.hex

    # Font recommendation
    if report.fonts:
        report.recommended_font = report.fonts.minor_latin or "Segoe UI"

    # Layout mapping — match template layouts to scenario layout types
    mapping = {}
    for sl in report.layouts:
        name = sl.name.lower()
        if sl.is_dark:
            continue  # Skip dark variants, prefer light

        if name.strip() in ("blank", "blank layout"):
            mapping["blank"] = sl.name
        elif name.strip() in ("title", "title slide") and "image" not in name:
            if "cover" not in mapping:
                mapping["cover"] = sl.name
        elif "title slide" in name and "image" in name:
            if "cover" not in mapping:
                mapping["cover"] = sl.name
        elif name.strip() in ("title and content", "object", "two_objects"):
            if "full_content" not in mapping:
                mapping["full_content"] = sl.name
        elif name.strip() in ("title only", "title_only"):
            mapping["title_only"] = sl.name
        elif name.strip() in ("section header", "section_header"):
            mapping["section_divider"] = sl.name
        elif name.strip() in ("two objects", "two_objects", "two_objects_with_text"):
            if "two_column" not in mapping:
                mapping["two_column"] = sl.name
        elif "content" in name and "picture" in name and "right" in name:
            mapping["screenshot_right"] = sl.name
        elif "picture" in name and "left" in name and "content" in name:
            mapping["screenshot_left"] = sl.name
        elif "agenda" in name:
            mapping["agenda"] = sl.name
        elif "end slide" in name or "logo slide" in name:
            mapping["closing"] = sl.name
        elif "question" in name and "answer" in name:
            mapping["qa"] = sl.name
        elif "testimonial" in name:
            mapping["testimonial"] = sl.name
        elif ("four" in name or "three" in name or "five" in name) and "grid" in name:
            if "grid" not in mapping:
                mapping["grid"] = sl.name
        elif "collage" in name:
            if "image_collage" not in mapping:
                mapping["image_collage"] = sl.name
        elif "infographic" in name:
            if "infographic" not in mapping:
                mapping["infographic"] = sl.name

    report.layout_mapping = mapping


# ---------------------------------------------------------------------------
# Output formatting
# ---------------------------------------------------------------------------

def format_report(report: TemplateReport) -> str:
    """Format report as human-readable text for Claude context."""
    lines = []
    lines.append(f"# Template Analysis: {os.path.basename(report.file_path)}")
    lines.append(f"")
    lines.append(f"## File Info")
    lines.append(f"- Path: `{report.file_path}`")
    lines.append(f"- Size: {report.file_size_kb:.1f} KB")
    lines.append(f"- Dimensions: {report.slide_width} x {report.slide_height} Emu ({report.aspect_ratio})")
    lines.append(f"- Existing slides: {report.existing_slides}")

    lines.append(f"")
    lines.append(f"## Theme")
    lines.append(f"- Name: {report.theme_name}")
    lines.append(f"- Color scheme: {report.color_scheme_name}")

    if report.fonts:
        lines.append(f"- Heading font: {report.fonts.major_latin}")
        lines.append(f"- Body font: {report.fonts.minor_latin}")
        if report.fonts.major_ea:
            lines.append(f"- East Asian heading: {report.fonts.major_ea}")
        if report.fonts.minor_ea:
            lines.append(f"- East Asian body: {report.fonts.minor_ea}")

    lines.append(f"")
    lines.append(f"## Theme Colors")
    lines.append(f"| Role | HEX | Semantic |")
    lines.append(f"|------|-----|----------|")
    for tc in report.colors:
        semantic = ROLE_MAP.get(tc.role, "")
        lines.append(f"| {tc.role} | `{tc.hex}` | {semantic} |")

    if report.design_colors:
        lines.append(f"")
        lines.append(f"## Design Colors (from slide shapes)")
        lines.append(f"Custom theme: {'Yes' if report.is_custom_theme else 'No (default Office → colors extracted from shapes)'}")
        lines.append(f"| HEX | Count |")
        lines.append(f"|-----|-------|")
        for hex_val, count in report.design_colors:
            lines.append(f"| `{hex_val}` | {count} |")

    lines.append(f"")
    lines.append(f"## Recommended Color Map (for scenario frontmatter)")
    lines.append(f"```yaml")
    lines.append(f"colors:")
    for semantic, hex_val in report.recommended_colors.items():
        lines.append(f'  {semantic}: "{hex_val}"')
    lines.append(f"```")

    lines.append(f"")
    lines.append(f"## Slide Layouts ({len(report.layouts)} total)")
    lines.append(f"| # | Name | Title | Content | Picture | Dark |")
    lines.append(f"|---|------|-------|---------|---------|------|")
    for sl in report.layouts:
        t = "Y" if sl.has_title else ""
        c = "Y" if sl.has_content else ""
        p = "Y" if sl.has_picture else ""
        d = "Y" if sl.is_dark else ""
        lines.append(f"| {sl.index} | {sl.name} | {t} | {c} | {p} | {d} |")

    lines.append(f"")
    lines.append(f"## Layout Mapping (template -> scenario)")
    lines.append(f"| Scenario Type | Template Layout |")
    lines.append(f"|---------------|-----------------|")
    for scenario_type, template_name in report.layout_mapping.items():
        lines.append(f"| `{scenario_type}` | {template_name} |")

    if report.recommended_font:
        lines.append(f"")
        lines.append(f"## Recommended Font")
        lines.append(f'`font: "{report.recommended_font}"`')

    return "\n".join(lines)


def to_json(report: TemplateReport) -> str:
    """Serialize report as JSON."""
    return json.dumps(asdict(report), indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 pptx_inspector.py <template.pptx> [--json]")
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.exists(path):
        print(f"File not found: {path}")
        sys.exit(1)

    report = inspect_template(path)

    if "--json" in sys.argv:
        print(to_json(report))
    else:
        print(format_report(report))


if __name__ == "__main__":
    main()
