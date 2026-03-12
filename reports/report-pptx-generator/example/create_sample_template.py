#!/usr/bin/env python3
"""Create a modern sample template.pptx with embedded theme colors and fonts.

This generates a professional starter template for use with report-pptx-generator,
so users don't need to bring their own corporate template.

Usage:
    python3 create_sample_template.py

Output:
    sample-template.pptx (in the same directory)
"""

from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(SCRIPT_DIR, "sample-template.pptx")

# Modern color palette — "Slate & Teal"
COLORS = {
    "dk1":      "1E293B",  # Slate 800 — primary dark text
    "lt1":      "FFFFFF",  # White
    "dk2":      "0F766E",  # Teal 700 — primary brand
    "lt2":      "E2E8F0",  # Slate 200 — light background
    "accent1":  "0D9488",  # Teal 600
    "accent2":  "64748B",  # Slate 500
    "accent3":  "F59E0B",  # Amber 500 — accent
    "accent4":  "6366F1",  # Indigo 500
    "accent5":  "1E40AF",  # Blue 800 — secondary
    "accent6":  "E11D48",  # Rose 600
    "hlink":    "0284C7",  # Sky 600
    "folHlink": "94A3B8",  # Slate 400
}

MAJOR_FONT = "Inter"       # Modern heading font (falls back to system sans)
MINOR_FONT = "Inter"       # Modern body font
FALLBACK_FONT = "Segoe UI"  # Widely available fallback


def inject_theme(prs: Presentation):
    """Replace the default theme with our custom color and font scheme."""
    master = prs.slide_masters[0]

    # Find theme part
    theme_part = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            theme_part = rel.target_part
            break

    if theme_part is None:
        print("Warning: No theme part found, skipping theme injection")
        return

    theme_xml = etree.fromstring(theme_part.blob)
    theme_xml.set("name", "PPTX Skills — Slate & Teal")

    # --- Replace color scheme ---
    for clrScheme in theme_xml.iter(qn("a:clrScheme")):
        clrScheme.set("name", "Slate_Teal")
        for child in list(clrScheme):
            role = child.tag.split("}")[-1]
            if role in COLORS:
                # Clear existing children
                for sub in list(child):
                    child.remove(sub)
                # Add srgbClr
                srgb = etree.SubElement(child, qn("a:srgbClr"))
                srgb.set("val", COLORS[role])

    # --- Replace font scheme ---
    for fontScheme in theme_xml.iter(qn("a:fontScheme")):
        fontScheme.set("name", "Modern_Inter")
        for fontType in ["majorFont", "minorFont"]:
            fe = fontScheme.find(qn(f"a:{fontType}"))
            if fe is not None:
                latin = fe.find(qn("a:latin"))
                if latin is not None:
                    font = MAJOR_FONT if fontType == "majorFont" else MINOR_FONT
                    latin.set("typeface", font)
                ea = fe.find(qn("a:ea"))
                if ea is not None:
                    ea.set("typeface", "")

    # Write back
    theme_part._blob = etree.tostring(theme_xml, xml_declaration=True, encoding="UTF-8", standalone=True)


def add_sample_slides(prs: Presentation):
    """Add a few sample slides to show the template capabilities."""
    blank = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank = layout
            break
    if blank is None:
        blank = prs.slide_layouts[0]

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height

    PRIMARY = RGBColor(0x0F, 0x76, 0x6E)
    ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
    DARK = RGBColor(0x1E, 0x29, 0x3B)
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT = RGBColor(0x33, 0x41, 0x55)
    TEXT_LIGHT = RGBColor(0x64, 0x74, 0x8B)

    FONT = FALLBACK_FONT
    FONT_BOLD = FALLBACK_FONT

    def _shape(slide, l, t, w, h, fill=None, line=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.line.fill.background()
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        else:
            s.fill.background()
        if line:
            s.line.color.rgb = line
            s.line.fill.solid()
        return s

    def _text(slide, txt, l, t, w, h, sz=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        return box

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(blank)
    _shape(slide, Emu(0), Emu(0), Emu(5200000), SLIDE_H, fill=DARK)
    _shape(slide, Emu(0), Emu(0), Emu(45000), Emu(4000000), fill=ACCENT)

    _text(slide, "PPTX SKILLS", Emu(450000), Emu(800000), Emu(4300000), Emu(200000),
          sz=11, bold=True, color=ACCENT, font=FONT_BOLD)
    _text(slide, "Sample\nTemplate", Emu(450000), Emu(1200000), Emu(4300000), Emu(1200000),
          sz=40, bold=True, color=WHITE, font=FONT_BOLD)
    _shape(slide, Emu(450000), Emu(2500000), Emu(500000), Emu(30000), fill=ACCENT)
    _text(slide, "A modern, editable presentation\ngenerated with python-pptx",
          Emu(450000), Emu(2650000), Emu(4300000), Emu(500000),
          sz=14, color=RGBColor(0x94, 0xA3, 0xB8))
    _text(slide, "March 2026", Emu(450000), Emu(5800000), Emu(2000000), Emu(200000),
          sz=10, color=RGBColor(0x94, 0xA3, 0xB8))

    # Right side — feature cards
    _text(slide, "WHAT THIS DEMONSTRATES", Emu(5600000), Emu(500000), Emu(5500000), Emu(250000),
          sz=10, bold=True, color=TEXT_LIGHT, font=FONT_BOLD)

    cards = [
        ("Theme Colors", "12 embedded theme colors extracted by pptx_inspector.py", PRIMARY),
        ("Font Scheme", "Inter / Segoe UI with heading and body variants", RGBColor(0x1E, 0x40, 0xAF)),
        ("Slide Layouts", "Blank layout for maximum control with custom shapes", RGBColor(0x63, 0x66, 0xF1)),
        ("Editable Elements", "Every shape, text box, and table is editable in PowerPoint", ACCENT),
    ]
    for i, (title, desc, accent) in enumerate(cards):
        y = Emu(900000 + i * 1300000)
        from pptx.enum.shapes import MSO_SHAPE as S
        card = slide.shapes.add_shape(S.ROUNDED_RECTANGLE, Emu(5600000), y, Emu(6100000), Emu(1100000))
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        _shape(slide, Emu(5600000), y, Emu(40000), Emu(1100000), fill=accent)
        _text(slide, title, Emu(5850000), y + Emu(120000), Emu(5600000), Emu(300000),
              sz=15, bold=True, color=DARK, font=FONT_BOLD)
        _text(slide, desc, Emu(5850000), y + Emu(480000), Emu(5600000), Emu(450000),
              sz=11, color=TEXT_LIGHT)

    # --- Slide 2: Content Demo ---
    slide = prs.slides.add_slide(blank)
    _shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(54000), fill=PRIMARY)
    _text(slide, "Content Layout Demo", Emu(500000), Emu(280000), Emu(10000000), Emu(500000),
          sz=24, bold=True, color=DARK, font=FONT_BOLD)
    _shape(slide, Emu(500000), Emu(720000), Emu(700000), Emu(36000), fill=ACCENT)

    # Left column
    _text(slide, "Left Column", Emu(500000), Emu(900000), Emu(5500000), Emu(300000),
          sz=16, bold=True, color=PRIMARY, font=FONT_BOLD)

    t = slide.shapes.add_table(4, 3, Emu(500000), Emu(1250000), Emu(5500000), Emu(1500000)).table
    t.columns[0].width = Emu(1800000)
    t.columns[1].width = Emu(2200000)
    t.columns[2].width = Emu(1500000)
    for i, h in enumerate(["Feature", "Description", "Status"]):
        c = t.cell(0, i)
        c.text = h
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = PRIMARY

    rows = [
        ["Theme Colors", "Extracted from template", "Ready"],
        ["Font Scheme", "Heading + Body fonts", "Ready"],
        ["Layouts", "8 scenario types supported", "Ready"],
    ]
    for r, row in enumerate(rows):
        bg = LIGHT_BG if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cl = t.cell(r + 1, c)
            cl.text = val
            p = cl.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT
            p.font.name = FONT
            cl.fill.solid()
            cl.fill.fore_color.rgb = bg

    # Right column
    _text(slide, "Right Column", Emu(6200000), Emu(900000), Emu(5500000), Emu(300000),
          sz=16, bold=True, color=PRIMARY, font=FONT_BOLD)

    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(6200000), Emu(1250000), Emu(5500000), Emu(600000))
    callout.line.fill.background()
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    tf = callout.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_top = Emu(80000)
    p = tf.paragraphs[0]
    p.text = "This template is generated by python-pptx.\nAll elements are fully editable in PowerPoint."
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
    p.font.name = FONT

    # Color swatches
    _text(slide, "Theme Color Palette", Emu(6200000), Emu(2050000), Emu(5500000), Emu(300000),
          sz=12, bold=True, color=TEXT, font=FONT_BOLD)

    swatch_colors = [
        ("#0F766E", "Primary"), ("#F59E0B", "Accent"), ("#1E40AF", "Secondary"),
        ("#6366F1", "Indigo"), ("#E11D48", "Rose"), ("#64748B", "Slate"),
    ]
    for i, (hex_c, label) in enumerate(swatch_colors):
        x = Emu(6200000 + (i % 3) * 1900000)
        y = Emu(2350000 + (i // 3) * 700000)
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        _shape(slide, x, y, Emu(350000), Emu(350000), fill=RGBColor(r, g, b))
        _text(slide, f"{label}\n{hex_c}", x + Emu(430000), y + Emu(30000), Emu(1300000), Emu(300000),
              sz=9, color=TEXT_LIGHT)

    # Footer
    _shape(slide, Emu(500000), Emu(6400000), Emu(11200000), Emu(9000), fill=RGBColor(0xE2, 0xE8, 0xF0))
    _text(slide, "2", Emu(11500000), Emu(6500000), Emu(400000), Emu(250000),
          sz=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

    # --- Slide 3: Section Divider ---
    slide = prs.slides.add_slide(blank)
    _shape(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=PRIMARY)
    _shape(slide, Emu(0), Emu(3250000), SLIDE_W, Emu(5000), fill=ACCENT)
    _text(slide, "Section Title", Emu(1000000), Emu(2500000), Emu(10000000), Emu(600000),
          sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
    _text(slide, "This is a section divider slide", Emu(1000000), Emu(3450000), Emu(10000000), Emu(400000),
          sz=16, color=RGBColor(0x99, 0xF6, 0xE4), align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    inject_theme(prs)
    add_sample_slides(prs)

    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT)
    print(f"Created: {OUTPUT}")
    print(f"Size: {size / 1024:.1f} KB")
    print(f"Slides: {len(prs.slides)}")
    print(f"Theme: PPTX Skills — Slate & Teal")
    print(f"Colors: Primary #0F766E, Accent #F59E0B, Secondary #1E40AF")
    print(f"Font: {FALLBACK_FONT} (Inter in theme, Segoe UI fallback)")
    print(f"\nRun inspector to verify:")
    print(f"  python3 ../pptx_inspector.py {OUTPUT}")


if __name__ == "__main__":
    main()
